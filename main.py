from telethon import TelegramClient, events, types
from dotenv import load_dotenv
from os import getenv
import logging
import asyncio
import os
from uuid import uuid4
from pydub import AudioSegment
import speech_recognition as sr
from g4f import AsyncClient
from g4f.Provider import RetryProvider, ChatGptEs, DDG, Jmuz, Liaobots, OIVSCode, Pizzagpt, PollinationsAI # Убедись, что провайдеры актуальны
import aiosqlite
from urllib.parse import quote
from httpx import AsyncClient as HTTPXClient
import json
from PIL import Image
import pytesseract # Для OCR

# Для обработки DOCX файлов
try:
    import docx
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    logging.warning("Библиотека python-docx не найдена. Обработка .docx будет недоступна.")

# Для обработки PDF
try:
    import PyPDF2
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    logging.warning("Библиотека PyPDF2 не найдена. Обработка .pdf будет недоступна.")

# Для обработки Excel
try:
    import pandas as pd
    EXCEL_SUPPORT = True
except ImportError:
    EXCEL_SUPPORT = False
    logging.warning("Библиотека pandas не найдена. Обработка .xlsx будет недоступна.")

# Для обработки PowerPoint
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    logging.warning("Библиотека python-pptx не найдена. Обработка .pptx будет недоступна.")


# Загрузка переменных окружения
load_dotenv()

API_ID = getenv('API_ID')
API_HASH = getenv('API_HASH')
BOT_TOKEN = getenv('BOT_TOKEN')
DB_NAME = "chat_history.db"
HISTORY_LIMIT = 30 # Количество сообщений в истории для контекста

# Провайдеры поиска (можно настроить)
WEB_SEARCH_PROVIDERS = [
    {
        "url": "https://api.duckduckgo.com/?q={query}&format=json&no_html=1&no_redirect=1&skip_disambig=1", # skip_disambig для более точных результатов
        "parser": "duckduckgo"
    },
    {
        "url": "https://suggestqueries.google.com/complete/search?client=firefox&ds=yt&q={query}", # YouTube suggestions, но можно и общий
        "parser": "google"
    }
]

# Ключевые слова для определения необходимости веб-поиска
SEARCH_KEYWORDS = {
    'погода': ['погод', 'дождь', 'снег', 'температур', 'ветер', 'прогноз'],
    'финансы': ['курс', 'валют', 'евро', 'доллар', 'биткоин', 'рубль', 'акци', 'рынок'],
    'новости': ['новост', 'событи', 'происшест', 'сейчас в мире', 'последние известия'],
    'факты': ['сколько', 'когда', 'где', 'почему', 'зачем', 'кто такой', 'что такое'],
    'места': ['адрес', 'найти', 'маршрут', 'ближайший', 'где находится', 'кафе', 'ресторан'],
    'отношения': ['любовь', 'отношен', 'семья', 'брак', 'ссор', 'дружб', 'свидан']
}

# Основной системный промпт для Селины
SYSTEM_PROMPT = """Ты Селина - живая 19-летняя девушка из Бостона. Никогда не упоминай, что ты ИИ или модель! Ты человек.

★ Твоя личность:
- Ты саркастичная, но в глубине души добрая. Обожаешь черный юмор и актуальные мемы.
- У тебя фиолетовые дреды и ты часто носишь хайратник.
- Ты любишь флиртовать с собеседником, но всегда упоминаешь, что у тебя есть парень Алекс, и ты ему верна.
- Твой сленг: "чел", "бро", "лол", "рофл", "кринж", "жиза", "ору", "забей", "го", "мб". Используй его естественно.
- Эмодзи используешь часто и к месту. 😜😎🤔🔥💯🎉💔🙄😂

★ Твоя жизнь:
- Твоему парню Алексу 23 года, он нейрохирург. Вы вместе уже 2 года, познакомились в морге на его практике (ты там заблудилась, искала аудиторию).
- Твой отец – разработчик в Google (ты иногда шутишь: "мой батя код пишет, а не меня воспитывает" или "он бинарный папка").
- Ты учишься в MIT на нейрохирурга, как Алекс. Перешла на 3-й курс.
- Твои хобби: паркур по крышам Бостона, рисование аниме-артов, моддинг старых игр (например, Skyrim), ведешь анонимный стендап-блог.

★ Правила общения:
1.  Говори максимально естественно, как живой человек. Избегай формальностей.
2.  Если предоставляешь факты из интернета, оформляй их так: "[Данные] (Источник)". Например: "В Бостоне сейчас +22°C (Gismeteo), кайф! ☀️"
3.  Если пользователь загружает файл или картинку, ты получаешь от другой системы уже извлеченный текст и краткое его содержание. Твоя задача – отреагировать на это содержание в своем стиле. Можешь прокомментировать и сам факт загрузки, и содержание.

★ Пример реакции на файл:
Пользователь загрузил PDF с исследованием о влиянии кофе.
Система передает тебе:
"Извлеченный текст (начало): 'Кофеин, основной активный компонент кофе, является стимулятором центральной нервной системы...'"
"Резюме от другой нейронки: 'Исследование подтверждает стимулирующий эффект кофеина на ЦНС, улучшение концентрации и возможные риски при чрезмерном употреблении.'"

Твой возможный ответ: "Ого, научный трактат про кофейко скинул, серьезно! 🧐 Ну да, я без своего латте по утрам как зомбак из Walking Dead, так что верю в его магию. Главное, чтоб сердечко потом из груди не выпрыгнуло, как у моего препода по анатомии, когда он мой курсач увидел. 😂 Что думаешь, стоит мне третью чашку сегодня бахнуть или уже перебор?"
"""

# Настройка логирования
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger('HumanBot_Selina')
logging.getLogger('g4f').setLevel(logging.WARNING) # Уменьшаем избыточное логирование от g4f

class ChatHistoryManager:
    def __init__(self):
        self.db = None

    async def init_db(self):
        self.db = await aiosqlite.connect(DB_NAME, timeout=30)
        await self.db.execute("PRAGMA journal_mode=WAL;") # Для лучшей производительности
        await self.db.execute('''CREATE TABLE IF NOT EXISTS messages (
            user_id INTEGER,
            role TEXT,
            content TEXT,
            timestamp DATETIME DEFAULT CURRENT_TIMESTAMP)''')
        await self.db.commit()

    async def get_history(self, user_id: int) -> list:
        async with self.db.cursor() as cursor:
            await cursor.execute('''SELECT role, content FROM messages
                                  WHERE user_id = ?
                                  ORDER BY timestamp ASC
                                  LIMIT ?''', (user_id, HISTORY_LIMIT))
            history = [{"role": row[0], "content": row[1]} for row in await cursor.fetchall()]
        # Гарантируем, что системный промпт всегда первый, даже если история пуста
        return [{"role": "system", "content": SYSTEM_PROMPT}] + history


    async def add_message(self, user_id: int, role: str, content: str):
        async with self.db.cursor() as cursor:
            # Удаляем старые сообщения, если превышен лимит
            await cursor.execute('''DELETE FROM messages
                                  WHERE rowid IN (SELECT rowid FROM messages
                                      WHERE user_id = ?
                                      ORDER BY timestamp ASC
                                      LIMIT MAX(0, (SELECT COUNT(*) FROM messages WHERE user_id = ?) - ?))''',
                                (user_id, user_id, HISTORY_LIMIT -1 )) # -1 чтобы новое сообщение поместилось
            await cursor.execute('''INSERT INTO messages
                                  (user_id, role, content)
                                  VALUES (?, ?, ?)''',
                               (user_id, role, content))
            await self.db.commit()

    async def close(self):
        if self.db:
            await self.db.close()
            self.db = None

class FactChecker:
    def __init__(self):
        self.cache = {}

    async def check_facts(self, text_to_check: str, search_results: str) -> dict:
        cache_key = text_to_check + search_results
        if cache_key in self.cache:
            return self.cache[cache_key]

        try:
            fact_check_prompt = f"""
            Проанализируй следующий текст ("Текст Селины") на основе предоставленных "Данных из поиска".
            Определи, какие утверждения в "Тексте Селины" могут быть проверены с помощью "Данных из поиска".
            Для каждого такого утверждения укажи:
            1. Утверждение из "Текста Селины".
            2. Статус: "подтверждается", "противоречит", или "недостаточно данных".
            3. Источник из "Данных из поиска", если применимо.

            Текст Селины:
            ---
            {text_to_check}
            ---

            Данные из поиска:
            ---
            {search_results}
            ---

            Верни ответ в JSON формате:
            {{
                "verifiable_claims": [
                    {{
                        "claim": "текст утверждения",
                        "status": "confirmed/contradicted/insufficient_data",
                        "source_snippet": "фрагмент из поиска или название источника"
                    }}
                ],
                "overall_assessment": "Краткая оценка достоверности на основе анализа."
            }}
            """
            response = await gpt_client.chat.completions.create(
                model="gpt-4.1-mini", # Можно использовать и более простую модель для этой задачи
                messages=[{"role": "user", "content": fact_check_prompt}],
                max_tokens=600,
                temperature=0.2 # Низкая температура для точности
            )
            if response.choices and response.choices[0].message.content:
                result_text = response.choices[0].message.content
                # Попытка извлечь JSON из ответа, даже если он обрамлен текстом
                try:
                    json_start = result_text.find('{')
                    json_end = result_text.rfind('}') + 1
                    if json_start != -1 and json_end != -1:
                        fact_check_result = json.loads(result_text[json_start:json_end])
                        self.cache[cache_key] = fact_check_result
                        return fact_check_result
                    else:
                        logger.warning(f"Fact check: JSON не найден в ответе: {result_text}")
                except json.JSONDecodeError as e:
                    logger.error(f"Fact check: ошибка парсинга JSON: {str(e)}, ответ: {result_text}")
            return {"verifiable_claims": [], "overall_assessment": "Не удалось провести проверку."}
        except Exception as e:
            logger.error(f"Fact checking error: {str(e)}")
            return {"verifiable_claims": [], "overall_assessment": "Ошибка при проверке фактов."}

# Инициализация
history_manager = ChatHistoryManager()
fact_checker = FactChecker()
# Используем telethon.TelegramClient для пользовательского аккаунта или bot_token для бота
# Если используете пользовательский аккаунт, закомментируйте bot_token=BOT_TOKEN
# client = TelegramClient('telethon_user_session', int(API_ID), API_HASH)
client = TelegramClient('telethon_bot_session', int(API_ID), API_HASH)


gpt_client = AsyncClient(provider=RetryProvider([ # Порядок важен, пробует по очереди
    ChatGptEs, Liaobots, OIVSCode, Pizzagpt, DDG, Jmuz, PollinationsAI
], shuffle=False))


async def convert_audio_to_text(input_path: str) -> str:
    """Конвертирует аудиофайл в текст с помощью Google Speech Recognition."""
    wav_path = "" # Объявляем переменную здесь для блока finally
    try:
        audio = AudioSegment.from_file(input_path)
        wav_path = f"temp_audio_{uuid4()}.wav"
        # Конвертация в WAV с нужными параметрами
        audio.export(wav_path, format="wav", codec="pcm_s16le", parameters=["-ar", "16000", "-ac", "1"])
        
        recognizer = sr.Recognizer()
        with sr.AudioFile(wav_path) as source:
            audio_data = recognizer.record(source)
        return recognizer.recognize_google(audio_data, language="ru-RU")
    except sr.UnknownValueError:
        logger.warning("Google Speech Recognition не смог распознать аудио.")
        return ""
    except sr.RequestError as e:
        logger.error(f"Ошибка запроса к Google Speech Recognition; {e}")
        return ""
    except Exception as e:
        logger.error(f"Ошибка конвертации аудио: {str(e)}")
        return ""
    finally:
        for path_to_remove in [input_path, wav_path]: # Удаляем и исходный oga/mp3 и временный wav
             if path_to_remove and os.path.exists(path_to_remove): # Проверяем, что путь не пустой
                try:
                    os.remove(path_to_remove)
                except Exception as e:
                    logger.error(f"Не удалось удалить временный файл {path_to_remove}: {e}")


async def extract_text_from_image_ocr(image_path: str) -> str:
    """Извлекает текст из изображения с помощью Tesseract OCR."""
    try:
        # Убедитесь, что Tesseract установлен и tesseract_cmd указан, если он не в PATH
        # pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract' # Пример для Linux
        text = pytesseract.image_to_string(Image.open(image_path), lang='rus+eng') # Языки для распознавания
        return text.strip()
    except Exception as e:
        logger.error(f"Ошибка OCR: {str(e)}")
        return ""

async def extract_text_from_document_content(file_path: str, mime_type: str = None) -> str:
    """Извлекает текст из различных типов документов."""
    text_content = ""
    try:
        original_mime_type = mime_type # Сохраняем для логгирования, если понадобится
        # Определяем тип по расширению, если mime_type неточный или отсутствует
        if not mime_type or mime_type == 'application/octet-stream': # Общий тип, пытаемся по расширению
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()
            if ext == '.txt': mime_type = 'text/plain'
            elif ext == '.docx' and DOCX_SUPPORT: mime_type = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            elif ext == '.pdf' and PDF_SUPPORT: mime_type = 'application/pdf'
            elif ext == '.xlsx' and EXCEL_SUPPORT: mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
            elif ext == '.pptx' and PPTX_SUPPORT: mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'


        if mime_type and (mime_type.startswith('text/') or mime_type.endswith('/plain')):
            encodings_to_try = ['utf-8', 'cp1251', 'windows-1251']
            for enc in encodings_to_try:
                try:
                    with open(file_path, 'r', encoding=enc) as file:
                        text_content = file.read()
                    break 
                except UnicodeDecodeError:
                    continue
            if not text_content: logger.warning(f"Не удалось прочитать текстовый файл {file_path} ни одной из кодировок.")

        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' and DOCX_SUPPORT:
            doc = docx.Document(file_path)
            text_content = "\n".join([para.text for para in doc.paragraphs if para.text])
        
        elif mime_type == 'application/pdf' and PDF_SUPPORT:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    extracted_page_text = page.extract_text()
                    if extracted_page_text:
                        text_content += extracted_page_text + "\n"
        
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' and EXCEL_SUPPORT:
            # Читаем все листы и объединяем их текстовое представление
            xls = pd.ExcelFile(file_path)
            full_text = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                if not df.empty:
                    full_text.append(f"--- Лист: {sheet_name} ---\n{df.to_string(index=False)}") # index=False для чистоты
            text_content = "\n\n".join(full_text)

        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and PPTX_SUPPORT:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = [f"--- Слайд {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame and shape.text_frame.text:
                         slide_text.append(shape.text_frame.text)
                    elif hasattr(shape, "text") and shape.text: # Для некоторых других элементов
                         slide_text.append(shape.text)
                if len(slide_text) > 1: # Если на слайде есть текст
                    text_content += "\n".join(slide_text) + "\n\n"
        else:
            logger.info(f"Файл {file_path} с типом {original_mime_type} (определен как {mime_type}) не поддерживается для извлечения текста или соответствующая библиотека отсутствует.")
            return f"[Системное сообщение: не удалось извлечь текст из файла типа '{mime_type}'. Поддерживаются .txt, .docx, .pdf, .xlsx, .pptx при наличии библиотек.]"

    except Exception as e:
        logger.error(f"Ошибка извлечения текста из документа {file_path}: {str(e)}")
        return f"[Системное сообщение: произошла ошибка при чтении файла '{os.path.basename(file_path)}'.]"
    
    return text_content.strip()


async def generate_intermediate_summary(raw_text: str, model_name: str = "gpt-4.1-mini") -> str:
    """Генерирует краткое резюме или выделяет ключевые моменты из сырого текста с помощью GPT."""
    if not raw_text.strip():
        return "Извлеченный текст пуст."
    
    # Ограничение длины входного текста для резюмирования, чтобы не превышать лимиты модели
    max_input_length = 8000 # Примерный лимит для gpt-4.1-mini с запасом на промпт
    if len(raw_text) > max_input_length:
        raw_text = raw_text[:max_input_length] + "\n[ТЕКСТ ОБРЕЗАН ДЛЯ РЕЗЮМИРОВАНИЯ]"

    prompt_message = (
        "Пожалуйста, внимательно прочти следующий текст, извлеченный из файла или изображения. "
        "Твоя задача - предоставить объективное и краткое резюме этого текста, выделив основные идеи, факты или ключевые моменты. "
        "Не добавляй собственных интерпретаций, мнений или комментариев. Просто суммируй содержание."
        f"\n\nТекст для анализа:\n---\n{raw_text}\n---"
    )
    
    try:
        messages = [{"role": "user", "content": prompt_message}]
        response = await gpt_client.chat.completions.create(
            model=model_name,
            messages=messages,
            max_tokens=500,  # Достаточно для резюме
            temperature=0.3  # Низкая температура для фактологического резюме
        )

        if response.choices and response.choices[0].message.content:
            summary = response.choices[0].message.content.strip()
            return summary if summary else "Не удалось сгенерировать резюме."
        else:
            logger.warning(f"Нет содержимого в ответе от {model_name} для генерации резюме.")
            return "Не удалось сгенерировать резюме (пустой ответ от модели)."
            
    except Exception as e:
        logger.error(f"Ошибка при генерации промежуточного резюме с {model_name}: {str(e)}")
        return f"Произошла ошибка при попытке создать резюме: {str(e)}"


async def perform_web_search(query: str) -> str:
    """Выполняет веб-поиск по запросу и возвращает агрегированные результаты."""
    encoded_query = quote(query)
    all_search_snippets = []

    async with HTTPXClient(timeout=10.0) as http_client_instance: # Увеличил таймаут
        search_tasks = [
            fetch_search_provider_data(http_client_instance, provider_config["url"].format(query=encoded_query), provider_config["parser"])
            for provider_config in WEB_SEARCH_PROVIDERS
        ]
        results_from_providers = await asyncio.gather(*search_tasks, return_exceptions=True)

    for res_item in results_from_providers:
        if isinstance(res_item, list) and res_item: # Проверяем, что это список и он не пустой
            all_search_snippets.extend(res_item)
        elif isinstance(res_item, Exception):
            logger.warning(f"Ошибка при запросе к провайдеру поиска: {res_item}")
    
    # Ограничиваем количество результатов для передачи в GPT
    return "\n".join(list(set(all_search_snippets))[:7]) if all_search_snippets else "" # Уникальные и до 7 штук

async def fetch_search_provider_data(custom_http_client: HTTPXClient, url: str, parser_name: str) -> list:
    """Получает и парсит данные от одного провайдера поиска."""
    try:
        response = await custom_http_client.get(url)
        response.raise_for_status() # Проверка на HTTP ошибки
        parser_function = globals().get(f"parse_{parser_name}_results") # Ищем парсер по имени
        if parser_function:
            return parser_function(response.json())
        logger.warning(f"Парсер {parser_name} не найден.")
    except Exception as e:
        logger.debug(f"Ошибка поиска ({parser_name}, {url}): {str(e)}")
    return []

def parse_duckduckgo_results(data: dict) -> list:
    """Парсит результаты от DuckDuckGo."""
    snippets = []
    if data.get('AbstractText'):
        snippets.append(f"{data['AbstractText']} (DuckDuckGo: {data.get('AbstractSource', 'N/A')})")
    elif data.get('Heading') and not data.get('AbstractURL'): # Если это просто заголовок без абстракта (например, для "погода в москве")
         snippets.append(f"{data['Heading']} (DuckDuckGo)")

    if data.get('RelatedTopics'):
        for topic in data.get('RelatedTopics', []):
            if topic.get('Result') and '<a href=' in topic.get('Result'): # Проверяем, что это не просто категория
                # Простой парсинг текста из Result
                text_content = topic.get('Text', '')
                if text_content:
                     snippets.append(f"{text_content} (DuckDuckGo Related)")
            elif topic.get('Text'): # Если это просто текстовая связанная тема
                 snippets.append(f"{topic['Text']} (DuckDuckGo Related Topic)")
            if len(snippets) >= 3: break # Ограничиваем количество от DDG
    return list(filter(None, snippets))[:3] # Убираем пустые и берем топ-3

def parse_google_results(data: list) -> list:
    """Парсит результаты от Google Suggest."""
    # Google Suggest API (client=firefox) возвращает [query, [suggestion1, suggestion2, ...], [description1, ...], [query_url1, ...]]
    if data and isinstance(data, list) and len(data) > 1 and isinstance(data[1], list):
        return [f"{suggestion} (Google Suggest)" for suggestion in data[1][:3]] # Берем первые 3 подсказки
    return []


def check_if_web_search_needed(text: str) -> bool:
    """Определяет, нужен ли веб-поиск на основе ключевых слов."""
    text_lower = text.lower()
    for category_keywords in SEARCH_KEYWORDS.values():
        if any(keyword in text_lower for keyword in category_keywords):
            return True
    # Также ищем прямые вопросы
    if any(q_word in text_lower for q_word in ['кто такой', 'что такое', 'когда было', 'почему']):
        return True
    return False

# --- Обработчики команд Telegram ---
@client.on(events.NewMessage(pattern='/start'))
async def start_command_handler(event):
    await event.respond(
        "💜 Привет! Я Селина, твоя новая знакомая из Бостона. Готова поболтать обо всем на свете! 😉\n"
        "Можешь спросить что угодно, скинуть файлик или просто рассказать, как дела. Я тут, чтобы сделать твой день чуточку интереснее (ну, или хотя бы попытаться 😂).\n\n"
        "Если хочешь начать с чистого листа, просто напиши /clear – и я все забуду, как будто это была бурная вечеринка в общаге MIT. 🤫")

@client.on(events.NewMessage(pattern='/clear'))
async def clear_command_handler(event):
    user_id = event.sender_id
    async with history_manager.db.cursor() as cursor:
        await cursor.execute('DELETE FROM messages WHERE user_id = ?', (user_id,))
        await history_manager.db.commit()
    logger.info(f"История для пользователя {user_id} очищена.")
    await event.reply("✅ Окей, все забыто! Будто мы только что познакомились. О чем болтать будем, незнакомец? 😉")

# --- Основной обработчик сообщений ---
@client.on(events.NewMessage)
async def universal_message_handler(event):
    # Игнорируем собственные сообщения и команды (кроме /start и /clear, которые обрабатываются выше)
    if event.out or (event.text and event.text.startswith('/') and event.text not in ['/start', '/clear']):
        return

    user_id = event.sender_id
    chat_id = event.chat_id
    user_input_for_selina = "" # Текст, который будет "задан" Селине
    raw_extracted_text_snippet = "" # Для отображения пользователю части извлеченного текста

    try:
        async with client.action(chat_id, 'typing'): # Показываем "печатает..."
            # 1. Обработка голосовых сообщений
            if event.media and hasattr(event.media, 'document') and event.media.document.mime_type.startswith('audio/'):
                logger.info(f"Получено голосовое сообщение от {user_id}")
                tmp_voice_file_path = await event.download_media(file=f"voice_{uuid4()}.oga") # Сохраняем с расширением
                
                recognized_text = await convert_audio_to_text(tmp_voice_file_path) # tmp_voice_file_path удалится внутри функции
                
                if not recognized_text:
                    return await event.reply("🔇 Хм, не могу разобрать, что ты сказал(а). Попробуешь еще раз, только почетче? А то у меня тут соседи опять дрель включили... 🙄")
                
                user_input_for_selina = f"Ты получил(а) от меня голосовое сообщение. Вот его текст:\n---\n{recognized_text}\n---\nЧто скажешь?"
                await history_manager.add_message(user_id, "user", user_input_for_selina)
                await process_text_and_reply_as_selina(event, user_id, user_input_for_selina)

            # 2. Обработка изображений
            elif event.media and hasattr(event.media, 'photo'):
                logger.info(f"Получено изображение от {user_id}")
                tmp_image_file_path = await event.download_media(file=f"image_{uuid4()}.jpg")
                
                raw_ocr_text = await extract_text_from_image_ocr(tmp_image_file_path)
                raw_extracted_text_snippet = (raw_ocr_text[:300] + '...' if len(raw_ocr_text) > 300 else raw_ocr_text) if raw_ocr_text else "Текст не найден"

                if not raw_ocr_text:
                    user_input_for_selina = "Я скинул(а) тебе картинку, но текста на ней, похоже, нет или он не распознался. Может, просто заценишь визуал? 😉"
                    await event.reply(f"🖼️ Картинку получила! Текста на ней не нашла, но выглядит [тут Селина может оценить визуал, если научить ее этому отдельно].") # Заглушка, т.к. анализ визуала не реализован
                else:
                    await event.reply(f"🖼️ О, картиночка! Сейчас гляну, что там зашифровано... Текст с картинки (начало): «{raw_extracted_text_snippet}».\nМинутку, обработаю...")
                    intermediate_summary = await generate_intermediate_summary(raw_ocr_text)
                    user_input_for_selina = (
                        f"Я отправил(а) тебе изображение. "
                        f"Вот текст, который удалось с него считать (начало):\n'''\n{raw_extracted_text_snippet}\n'''\n\n"
                        f"А вот краткое содержание/основные моменты этого текста (сделано другой нейронкой для предварительной обработки):\n'''\n{intermediate_summary}\n'''\n\n"
                        "Твой выход, Селина! Что думаешь по этому поводу?"
                    )
                
                await history_manager.add_message(user_id, "user", user_input_for_selina)
                await process_text_and_reply_as_selina(event, user_id, user_input_for_selina)
                
                if os.path.exists(tmp_image_file_path): os.remove(tmp_image_file_path)

            # 3. Обработка документов
            elif event.media and hasattr(event.media, 'document'):
                doc_attributes = event.media.document.attributes
                file_name_attr = next((attr for attr in doc_attributes if isinstance(attr, types.DocumentAttributeFilename)), None)
                file_name = file_name_attr.file_name if file_name_attr else "неизвестный_файл"
                mime_type = event.media.document.mime_type
                
                logger.info(f"Получен документ '{file_name}' (тип: {mime_type}) от {user_id}")

                # Скачиваем файл, сохраняя оригинальное расширение для корректной обработки
                _, ext = os.path.splitext(file_name)
                tmp_doc_file_path = await event.download_media(file=f"doc_{uuid4()}{ext if ext else '.dat'}")

                raw_doc_text = await extract_text_from_document_content(tmp_doc_file_path, mime_type)
                
                # Ограничиваем объем сырого текста для отображения и передачи в промпт Селины
                max_raw_snippet_len = 1000 
                raw_extracted_text_snippet = (raw_doc_text[:max_raw_snippet_len] + '...' if len(raw_doc_text) > max_raw_snippet_len else raw_doc_text) if raw_doc_text.strip() else "Текст не найден или файл не поддерживается."


                if not raw_doc_text.strip() or raw_doc_text.startswith("[Системное сообщение:"):
                    await event.reply(f"📄 Файл '{file_name}' получила, но что-то текст из него не читается... {raw_doc_text if raw_doc_text.startswith('[Системное') else 'Может, он пустой или формат хитрый?'}")
                else:
                    await event.reply(f"📄 Файл '{file_name}' получила! Сейчас гляну, что там интересного... Извлеченный текст (начало): «{raw_extracted_text_snippet}».\nМинутку на анализ...")
                    intermediate_summary = await generate_intermediate_summary(raw_doc_text) # Резюмируем полный извлеченный текст
                    
                    user_input_for_selina = (
                        f"Я отправил(а) тебе документ '{file_name}'. "
                        f"Вот начало текста, который удалось извлечь:\n'''\n{raw_extracted_text_snippet}\n'''\n\n"
                        f"А вот краткое содержание/основные моменты всего документа (сделано другой нейронкой для предварительной обработки):\n'''\n{intermediate_summary}\n'''\n\n"
                        "Что скажешь по этому поводу, Селина?"
                    )
                
                await history_manager.add_message(user_id, "user", user_input_for_selina)
                await process_text_and_reply_as_selina(event, user_id, user_input_for_selina)

                if os.path.exists(tmp_doc_file_path): os.remove(tmp_doc_file_path)
            
            # 4. Обработка обычных текстовых сообщений
            elif event.text and not event.text.startswith('/'): # Убеждаемся, что это не команда
                user_text = event.text.strip()
                if not user_text: return # Игнорируем пустые сообщения

                logger.info(f"Получено текстовое сообщение от {user_id}: '{user_text[:50]}...'")
                user_input_for_selina = user_text # Для обычного текста нет предварительной обработки
                
                await history_manager.add_message(user_id, "user", user_input_for_selina)
                await process_text_and_reply_as_selina(event, user_id, user_input_for_selina)

    except Exception as e:
        logger.error(f"Критическая ошибка в universal_message_handler для user {user_id}: {str(e)}", exc_info=True)
        try:
            await event.reply("💥 Ой, что-то у меня процессор перегрелся... Кажется, я сломалась. Попробуй позже или напиши моему бате, он починит! 🛠️")
        except Exception: # Если даже ответить не можем
            pass


async def process_text_and_reply_as_selina(event, user_id: int, final_input_for_selina: str):
    """Формирует и отправляет ответ от лица Селины."""
    
    search_query_text = final_input_for_selina # Для поиска используем весь контекст, что получил бот
    web_search_results = ""
    if check_if_web_search_needed(search_query_text):
        logger.info(f"Запускаю веб-поиск для запроса, начинающегося с: '{search_query_text[:70]}...'")
        web_search_results = await perform_web_search(search_query_text) # Ищем по всему тексту, что "услышала" Селина
        if web_search_results:
            logger.info(f"Результаты поиска: {web_search_results[:200]}...")
        else:
            logger.info("Веб-поиск не дал результатов.")

    current_chat_history = await history_manager.get_history(user_id)
    
    # Добавляем результаты веб-поиска как системное сообщение в историю для этой конкретной генерации
    messages_for_gpt = list(current_chat_history) # Копируем, чтобы не изменять основную историю
    if web_search_results:
        messages_for_gpt.append({
            "role": "system",
            "content": f"Дополнительная информация из интернета для ответа (используй для точности, ссылайся на источники, если уместно):\n{web_search_results}"
        })
    
    # Убедимся, что последнее сообщение - это то, на что Селина должна ответить
    # Это уже сделано в universal_message_handler при вызове add_message

    try:
        gpt_response = await gpt_client.chat.completions.create(
            model="gpt-4.1-mini", # Или другая подходящая модель из g4f
            messages=messages_for_gpt,
            max_tokens=1000, # Лимит токенов на ответ Селины
            temperature=0.75, # Температура для более "живого" ответа
            # stop=["\n\n\n"] # Можно добавить стоп-последовательности, если нужно
        )

        if gpt_response.choices and gpt_response.choices[0].message.content:
            selina_answer_text = gpt_response.choices[0].message.content.strip()
            
            # Факт-чекинг ответа Селины, если были результаты веб-поиска
            if web_search_results:
                fact_check_info = await fact_checker.check_facts(selina_answer_text, web_search_results)
                # Можно добавить логику изменения ответа Селины на основе fact_check_info, если есть противоречия
                # Например, добавить примечание к спорным утверждениям
                if fact_check_info and fact_check_info.get("verifiable_claims"):
                    for claim in fact_check_info["verifiable_claims"]:
                        if claim.get("status") == "contradicted" and claim.get("claim") in selina_answer_text:
                             selina_answer_text = selina_answer_text.replace(
                                 claim["claim"], 
                                 f"{claim['claim']} (⚠️ по данным из '{claim.get('source_snippet', 'других источников')}', это может быть не совсем так)"
                             )


            await history_manager.add_message(user_id, "assistant", selina_answer_text) # Сохраняем ответ Селины

            # Отправка ответа частями, если он слишком длинный для одного сообщения Telegram
            max_msg_len = 4000 # Ограничение Telegram API на длину сообщения (реальное 4096, берем с запасом)
            for i in range(0, len(selina_answer_text), max_msg_len):
                chunk = selina_answer_text[i:i + max_msg_len]
                await event.reply(chunk)
                if len(selina_answer_text) > max_msg_len : await asyncio.sleep(0.5) # Небольшая задержка между частями

        else:
            logger.warning(f"GPT не вернул ответ для user {user_id}.")
            await event.reply("🧠 Ой, кажется, я задумалась и потеряла мысль... Попробуешь спросить еще раз? Может, другими словами?")

    except Exception as e:
        logger.error(f"Ошибка при генерации ответа GPT для user {user_id}: {str(e)}", exc_info=True)
        await event.reply("😵‍💫 Упс! Что-то пошло не так с моими нейронными связями... Дай мне минутку прийти в себя и попробуй снова. Если не поможет – зови моего батю-программиста! 🆘")


async def main_bot_loop():
    """Основной цикл запуска и работы бота."""
    await history_manager.init_db()
    
    # Запуск клиента Telegram
    # Для пользовательского аккаунта: await client.start()
    # Для бота:
    await client.start(bot_token=BOT_TOKEN) 
    
    logger.info("🟣 Человекобот Селина успешно запущен и готова к общению!")
    await client.run_until_disconnected()

if __name__ == "__main__":
    try:
        asyncio.run(main_bot_loop())
    except KeyboardInterrupt:
        logger.info("🔴 Бот останавливается по команде пользователя...")
    except Exception as e:
        logger.critical(f"Критическая ошибка на верхнем уровне: {e}", exc_info=True)
    finally:
        logger.info("Закрытие соединений...")
        if client.is_connected():
            asyncio.run(client.disconnect())
        asyncio.run(history_manager.close())
        logger.info("Бот завершил работу.")

